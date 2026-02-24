"""
影幻智提 (VidSlide) - PPT 幻灯片智能提取工具 (多任务版)
======================================================
基于 Flask 的本地 Web 应用，提供可视化界面来提取、管理和打包 PPT 幻灯片。
支持同时对多个视频进行提取（最多 3 个并行标签页）。

使用方法：
    python app.py

依赖安装：
    pip install flask opencv-python numpy pillow python-pptx psutil

作者: PWO-CHINA
版本: v0.2.1
"""

import cv2
import gc
import numpy as np
import os
import sys
import shutil
import threading
import time
import uuid
import zipfile
import webbrowser
import socket
import traceback
from pathlib import Path

from flask import Flask, request, jsonify, send_file, send_from_directory, render_template
from PIL import Image

# ============================================================
#  无控制台模式兼容
# ============================================================
if sys.stdout is None:
    sys.stdout = open(os.devnull, 'w', encoding='utf-8')
if sys.stderr is None:
    sys.stderr = open(os.devnull, 'w', encoding='utf-8')

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️  未安装 python-pptx，PPTX 导出将不可用。安装命令: pip install python-pptx")

try:
    import psutil
    HAS_PSUTIL = True
except ImportError:
    HAS_PSUTIL = False
    print("⚠️  未安装 psutil，系统资源监控将不可用。安装命令: pip install psutil")


# ============================================================
#  PyInstaller / Nuitka 兼容：资源路径寻路
# ============================================================
def _is_frozen():
    """判断是否以打包后的 exe 运行（PyInstaller 或 Nuitka）"""
    return (getattr(sys, 'frozen', False)
            or hasattr(sys, '_MEIPASS')
            or '__compiled__' in globals())


def get_resource_path(relative_path):
    """获取打包后的资源文件路径"""
    if hasattr(sys, '_MEIPASS'):
        # PyInstaller 解压临时目录
        return os.path.join(sys._MEIPASS, relative_path)
    # Nuitka --include-data-dir 或源码模式：文件在 __file__ 旁边
    return os.path.join(os.path.abspath(os.path.dirname(__file__)), relative_path)


# ============================================================
#  配置
# ============================================================
TEMPLATE_DIR = get_resource_path('templates')

if _is_frozen():
    # PyInstaller / Nuitka exe: 用户文件放在 exe 所在目录
    BASE_DIR = os.path.dirname(os.path.abspath(sys.executable))
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# 多会话的根目录
SESSIONS_ROOT = os.path.join(BASE_DIR, '.vidslide_sessions')
MAX_SESSIONS = 3   # 最大并行标签页数

app = Flask(__name__, template_folder=TEMPLATE_DIR)


# ============================================================
#  全局错误处理 & CORS
# ============================================================
@app.after_request
def after_request(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
    response.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS'
    return response


@app.errorhandler(Exception)
def handle_error(e):
    traceback.print_exc()
    return jsonify(
        success=False,
        message=f'服务器内部错误: {str(e)}',
        error_type=type(e).__name__,
        hint='如果问题持续出现，请前往 https://github.com/PWO-CHINA/VidSlide/issues 提交 Issue，'
             '并附上此错误信息的截图。'
    ), 500


# ============================================================
#  多会话状态管理 (线程安全)
# ============================================================
_sessions_lock = threading.Lock()
_sessions = {}  # session_id -> session dict

# 资源告警阈值
CPU_WARN_THRESHOLD = 90
MEMORY_WARN_THRESHOLD = 85
DISK_WARN_THRESHOLD_MB = 500

# ── 后台 CPU 采样（避免 psutil.cpu_percent 阻塞请求线程）──
_cpu_cache = {'percent': 0.0}

def _cpu_sampler_loop():
    """后台线程：每 2 秒采样一次 CPU，结果写入 _cpu_cache"""
    if not HAS_PSUTIL:
        return
    psutil.cpu_percent(interval=1)  # 首次初始化
    while True:
        try:
            _cpu_cache['percent'] = psutil.cpu_percent(interval=0)
        except Exception:
            pass
        time.sleep(2)

_cpu_sampler_thread = threading.Thread(target=_cpu_sampler_loop, daemon=True)
_cpu_sampler_thread.start()


def _create_session():
    """创建一个新会话"""
    sid = uuid.uuid4().hex[:8]
    cache_dir = os.path.join(SESSIONS_ROOT, sid, 'cache')
    pkg_dir = os.path.join(SESSIONS_ROOT, sid, 'packages')
    os.makedirs(cache_dir, exist_ok=True)
    os.makedirs(pkg_dir, exist_ok=True)

    session = {
        'id': sid,
        'created_at': time.time(),
        'status': 'idle',
        'progress': 0,
        'message': '',
        'saved_count': 0,
        'video_path': '',
        'video_name': '',
        'cancel_flag': False,
        'eta_seconds': -1,
        'elapsed_seconds': 0,
        'cache_dir': cache_dir,
        'pkg_dir': pkg_dir,
        'lock': threading.Lock(),
    }
    with _sessions_lock:
        _sessions[sid] = session
    return sid


def _get_session(sid):
    with _sessions_lock:
        return _sessions.get(sid)


def _get_session_state(sid):
    sess = _get_session(sid)
    if not sess:
        return None
    with sess['lock']:
        return {k: v for k, v in sess.items() if k != 'lock'}


def _update_session(sid, **kw):
    sess = _get_session(sid)
    if not sess:
        return
    with sess['lock']:
        sess.update(kw)


def _delete_session(sid):
    with _sessions_lock:
        sess = _sessions.pop(sid, None)
    if sess:
        session_dir = os.path.join(SESSIONS_ROOT, sid)
        if os.path.exists(session_dir):
            shutil.rmtree(session_dir, ignore_errors=True)


def _get_all_sessions_summary():
    with _sessions_lock:
        sids = list(_sessions.keys())
    result = []
    for sid in sids:
        state = _get_session_state(sid)
        if state:
            result.append({
                'id': state['id'],
                'status': state['status'],
                'progress': state['progress'],
                'message': state['message'],
                'saved_count': state['saved_count'],
                'video_path': state['video_path'],
                'video_name': state['video_name'],
                'eta_seconds': state['eta_seconds'],
                'elapsed_seconds': state['elapsed_seconds'],
            })
    return result


def _count_running():
    with _sessions_lock:
        sids = list(_sessions.keys())
    count = 0
    for sid in sids:
        sess = _get_session(sid)
        if sess:
            with sess['lock']:
                if sess['status'] == 'running':
                    count += 1
    return count


# 心跳
_last_heartbeat = 0.0
_heartbeat_received = False
HEARTBEAT_TIMEOUT = 20


# ============================================================
#  路由 — 页面
# ============================================================
@app.route('/')
def index():
    return render_template('index.html')


# ============================================================
#  路由 — 会话管理
# ============================================================
@app.route('/api/sessions', methods=['GET'])
def list_sessions():
    return jsonify(success=True, sessions=_get_all_sessions_summary(), max_sessions=MAX_SESSIONS)


@app.route('/api/session/create', methods=['POST'])
def create_session():
    with _sessions_lock:
        current_count = len(_sessions)
    if current_count >= MAX_SESSIONS:
        return jsonify(success=False,
                       message=f'已达到最大标签页数量（{MAX_SESSIONS}个）。请先关闭不需要的标签页。')

    warning = _check_resource_warning()
    if warning:
        return jsonify(success=False, message=f'系统资源不足，无法新建标签页：{warning}')

    sid = _create_session()
    return jsonify(success=True, session_id=sid)


@app.route('/api/session/<sid>/close', methods=['POST'])
def close_session(sid):
    sess = _get_session(sid)
    if not sess:
        return jsonify(success=False, message='会话不存在')
    with sess['lock']:
        if sess['status'] == 'running':
            sess['cancel_flag'] = True
    time.sleep(0.3)
    _delete_session(sid)
    return jsonify(success=True)


# ============================================================
#  路由 — 选择视频
# ============================================================
@app.route('/api/select-video', methods=['POST'])
def select_video():
    try:
        try:
            import ctypes
            ctypes.windll.shcore.SetProcessDpiAwareness(2)
        except Exception:
            try:
                import ctypes
                ctypes.windll.user32.SetProcessDPIAware()
            except Exception:
                pass

        import tkinter as tk
        from tkinter import filedialog
        import queue

        result_queue = queue.Queue()

        def _pick():
            try:
                root = tk.Tk()
                root.withdraw()
                root.wm_attributes('-topmost', 1)
                root.focus_force()
                path = filedialog.askopenfilename(
                    title="请选择要提取的课程视频",
                    filetypes=[
                        ("视频文件", "*.mp4 *.avi *.mkv *.mov *.flv *.wmv *.webm"),
                        ("所有文件", "*.*"),
                    ],
                )
                root.destroy()
                result_queue.put(path or '')
            except Exception as e:
                print(f'[DEBUG] tkinter 弹窗异常: {e}')
                result_queue.put('')

        t = threading.Thread(target=_pick, daemon=True)
        t.start()
        t.join(timeout=120)

        path = result_queue.get_nowait() if not result_queue.empty() else ''

        if path:
            print(f'[DEBUG] 用户选择了视频: {path}')
            return jsonify(success=True, path=path)
        return jsonify(success=False, message='未选择文件')
    except Exception as e:
        print(f'[ERROR] select_video 异常: {e}')
        return jsonify(success=False, message=str(e))


# ============================================================
#  路由 — 开始提取（会话级）
# ============================================================
@app.route('/api/session/<sid>/extract', methods=['POST'])
def start_extraction(sid):
    sess = _get_session(sid)
    if not sess:
        return jsonify(success=False, message='会话不存在')

    with sess['lock']:
        if sess['status'] == 'running':
            return jsonify(success=False, message='该标签页正在提取中，请等待完成或取消')

    warning = _check_resource_warning()
    if warning:
        return jsonify(success=False, message=f'系统资源不足：{warning}\n请等待其他任务完成或关闭不需要的标签页。')

    data = request.json or {}
    video_path = data.get('video_path', '')
    print(f'[DEBUG][{sid}] 收到提取请求，视频路径: {repr(video_path)}')

    threshold = float(data.get('threshold', 5.0))
    enable_history = bool(data.get('enable_history', False))
    max_history = int(data.get('max_history', 5))
    use_roi = bool(data.get('use_roi', True))
    fast_mode = bool(data.get('fast_mode', True))

    if not video_path:
        return jsonify(success=False, message='未提供视频路径')

    if not os.path.exists(video_path):
        return jsonify(success=False, message=f'视频文件不存在: {video_path}',
                       hint='请检查文件是否已被移动或删除，然后重新选择视频。')

    # ── 视频文件预检测 ──
    try:
        _test_cap = cv2.VideoCapture(video_path)
        if not _test_cap.isOpened():
            _test_cap.release()
            return jsonify(success=False,
                           message='无法打开视频文件，可能文件已损坏或格式不支持。',
                           hint='建议：1) 检查文件是否完整下载；2) 尝试用播放器打开验证；'
                                '3) 如果是 m3u8 格式，请先用猫抓完整下载为 mp4。')
        _test_ok, _test_frame = _test_cap.read()
        _fourcc = int(_test_cap.get(cv2.CAP_PROP_FOURCC))
        _codec = ''.join([chr((_fourcc >> 8 * i) & 0xFF) for i in range(4)]) if _fourcc else 'N/A'
        _total = int(_test_cap.get(cv2.CAP_PROP_FRAME_COUNT))
        _fps = _test_cap.get(cv2.CAP_PROP_FPS) or 0
        _test_cap.release()
        if not _test_ok or _test_frame is None:
            return jsonify(success=False,
                           message=f'视频解码失败（编解码器: {_codec}）。',
                           hint='可能原因：1) 视频编码不被 OpenCV 支持；2) 文件不完整。'
                                '建议：尝试用 FFmpeg 转码为 mp4 后重试。')
        if _total < 10 or _fps <= 0:
            return jsonify(success=False,
                           message=f'视频信息异常：帧数={_total}，FPS={_fps:.1f}。',
                           hint='该文件可能不是有效的视频文件，或已严重损坏。')
        print(f'[DEBUG][{sid}] 视频预检通过: codec={_codec}, frames={_total}, fps={_fps:.1f}')
    except cv2.error as e:
        return jsonify(success=False,
                       message=f'OpenCV 视频检测出错: {str(e)}',
                       hint='可能是视频编码不兼容。建议用 FFmpeg 转码为 H.264 mp4 后重试。')
    except Exception as e:
        return jsonify(success=False,
                       message=f'视频文件预检测失败: {str(e)}',
                       hint='请确认文件路径正确且文件未被其他程序占用。')

    cache_dir = sess['cache_dir']
    if os.path.exists(cache_dir):
        shutil.rmtree(cache_dir)
    os.makedirs(cache_dir, exist_ok=True)

    video_name = Path(video_path).stem or '未命名视频'
    _update_session(sid,
        status='running', progress=0, message='正在初始化…',
        saved_count=0, video_path=video_path, video_name=video_name,
        cancel_flag=False, eta_seconds=-1, elapsed_seconds=0,
    )

    threading.Thread(
        target=_extract_worker,
        args=(sid, video_path, cache_dir, threshold, enable_history, max_history, use_roi, fast_mode),
        daemon=True,
    ).start()

    return jsonify(success=True)


# ============================================================
#  后台提取 Worker（多会话版）
# ============================================================
def _extract_worker(sid, video_path, output_dir, threshold, enable_history, max_history, use_roi, fast_mode=True):
    cap = None
    history_pool = None
    try:
        cap = cv2.VideoCapture(video_path)
        ok, prev_frame = cap.read()
        if not ok:
            _update_session(sid, status='error', message='无法读取视频文件')
            return

        total_frames = max(int(cap.get(cv2.CAP_PROP_FRAME_COUNT)), 1)
        fps = cap.get(cv2.CAP_PROP_FPS) or 30
        frame_step = max(1, int(fps))

        h, w = prev_frame.shape[:2]
        if use_roi:
            y1, y2 = int(h * 0.185), h
            x1, x2 = int(w * 0.208), w
        else:
            y1, y2 = 0, h
            x1, x2 = 0, w

        roi_w = x2 - x1
        COMPARE_WIDTH = 480
        if fast_mode and roi_w > COMPARE_WIDTH:
            _scale = COMPARE_WIDTH / roi_w
        else:
            _scale = 1.0

        def _to_gray(frame):
            roi = frame[y1:y2, x1:x2]
            if _scale < 1.0:
                roi = cv2.resize(roi, None, fx=_scale, fy=_scale,
                                 interpolation=cv2.INTER_AREA)
            return cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)

        prev_gray = _to_gray(prev_frame)
        history_pool = [prev_gray] if enable_history else None

        count = 0
        saved = 0
        _extract_start_time = time.time()
        _THROTTLE_INTERVAL = 0.008  # 每轮主循环让出 8ms CPU，降低峰值占用

        def _should_cancel():
            """快速检查取消标志"""
            s = _get_session(sid)
            if not s:
                return True
            with s['lock']:
                return s['cancel_flag']

        fp = os.path.join(output_dir, f"slide_{saved:04d}.jpg")
        cv2.imencode('.jpg', prev_frame, [cv2.IMWRITE_JPEG_QUALITY, 95])[1].tofile(fp)
        saved += 1
        _update_session(sid, saved_count=saved, message=f'已提取 {saved} 张')

        while True:
            if _should_cancel():
                _update_session(sid, status='cancelled', message=f'已取消，已保存 {saved} 张')
                return

            # — 节流：让出少量 CPU 给系统和其他线程 —
            time.sleep(_THROTTLE_INTERVAL)

            grabbed = True
            for _ in range(frame_step):
                count += 1
                if not cap.grab():
                    grabbed = False
                    break
            if not grabbed:
                break

            if _should_cancel():
                _update_session(sid, status='cancelled', message=f'已取消，已保存 {saved} 张')
                return

            ok, curr_frame = cap.retrieve()
            if not ok:
                break

            pct = min(99, int(count / total_frames * 100))
            elapsed = time.time() - _extract_start_time
            if pct > 2:
                eta = elapsed / pct * (100 - pct)
            else:
                eta = -1
            _update_session(sid, progress=pct, eta_seconds=round(eta, 1), elapsed_seconds=round(elapsed, 1))

            curr_gray = _to_gray(curr_frame)
            mean_diff = np.mean(cv2.absdiff(curr_gray, prev_gray))

            if mean_diff > threshold:
                check_step = max(1, int(fps * 0.5))
                stable = 0
                last_gray = curr_gray
                settled_frame = None
                settled_gray = None

                while True:
                    if _should_cancel():
                        break  # 跳出稳定帧检测，外层会处理取消
                    time.sleep(_THROTTLE_INTERVAL)  # 子循环也节流
                    s_grabbed = True
                    for _ in range(check_step):
                        count += 1
                        if not cap.grab():
                            s_grabbed = False
                            break
                    if not s_grabbed:
                        break
                    ret, tmp = cap.retrieve()
                    if not ret:
                        break
                    tmp_gray = _to_gray(tmp)
                    if np.mean(cv2.absdiff(tmp_gray, last_gray)) < 1.0:
                        stable += 1
                    else:
                        stable = 0
                    last_gray = tmp_gray
                    if stable >= 2:
                        settled_frame = tmp
                        settled_gray = tmp_gray
                        break

                # 稳定帧检测后再检查一次取消
                if _should_cancel():
                    _update_session(sid, status='cancelled', message=f'已取消，已保存 {saved} 张')
                    return

                if settled_gray is not None:
                    final_diff = np.mean(cv2.absdiff(settled_gray, prev_gray))
                    dup = False
                    if enable_history and history_pool:
                        for pg in history_pool:
                            if np.mean(cv2.absdiff(settled_gray, pg)) <= threshold:
                                dup = True
                                break
                    elif final_diff <= threshold:
                        dup = True

                    if not dup and final_diff > threshold:
                        fp = os.path.join(output_dir, f"slide_{saved:04d}.jpg")
                        cv2.imencode('.jpg', settled_frame, [cv2.IMWRITE_JPEG_QUALITY, 95])[1].tofile(fp)
                        saved += 1
                        _update_session(sid, saved_count=saved, message=f'已提取 {saved} 张')
                        prev_gray = settled_gray
                        if enable_history:
                            history_pool.append(settled_gray)
                            if len(history_pool) > max_history:
                                history_pool.pop(0)
                    else:
                        prev_gray = settled_gray

        elapsed_total = round(time.time() - _extract_start_time, 1)
        _update_session(sid, status='done', progress=100, eta_seconds=0, elapsed_seconds=elapsed_total,
               message=f'提取完成！共 {saved} 张幻灯片，耗时 {int(elapsed_total)}s')
    except Exception as e:
        error_detail = traceback.format_exc()
        print(f"！！！[{sid}] 发生严重错误！！！\n{error_detail}")
        # 为用户提供可操作的错误信息
        err_msg = str(e)
        if 'memory' in err_msg.lower() or 'MemoryError' in type(e).__name__:
            hint = '内存不足，请关闭其他标签页或程序后重试。'
        elif 'permission' in err_msg.lower() or 'access' in err_msg.lower():
            hint = '文件权限被拒绝，请检查文件是否正在被其他程序使用。'
        elif isinstance(e, cv2.error):
            hint = '视频处理出错，建议用 FFmpeg 转码后重试。'
        else:
            hint = '请截图此错误并前往 GitHub Issues 反馈。'
        _update_session(sid, status='error',
                        message=f'提取出错: {err_msg}\n💡 {hint}')
    finally:
        # ── 确保释放所有重量级资源 ──
        if cap is not None:
            try:
                cap.release()
            except Exception:
                pass
        cap = None
        history_pool = None
        # 立即触发垃圾回收，释放大量 numpy 数组占用的内存
        gc.collect()


# ============================================================
#  路由 — 进度 / 取消
# ============================================================
@app.route('/api/session/<sid>/progress')
def session_progress(sid):
    state = _get_session_state(sid)
    if not state:
        return jsonify(success=False, message='会话不存在'), 404
    return jsonify(state)


@app.route('/api/session/<sid>/cancel', methods=['POST'])
def session_cancel(sid):
    _update_session(sid, cancel_flag=True)
    return jsonify(success=True)


# ============================================================
#  路由 — 图片列表 / 提供图片
# ============================================================
@app.route('/api/session/<sid>/images')
def session_list_images(sid):
    sess = _get_session(sid)
    if not sess:
        return jsonify(images=[])
    cache_dir = sess['cache_dir']
    if not os.path.exists(cache_dir):
        return jsonify(images=[])
    imgs = sorted(
        f for f in os.listdir(cache_dir)
        if f.lower().endswith(('.jpg', '.jpeg', '.png'))
    )
    return jsonify(images=imgs)


@app.route('/api/session/<sid>/image/<path:filename>')
def session_serve_image(sid, filename):
    sess = _get_session(sid)
    if not sess:
        return jsonify(success=False, message='会话不存在'), 404
    resp = send_from_directory(sess['cache_dir'], filename)
    resp.headers['Cache-Control'] = 'no-store'
    return resp


# ============================================================
#  路由 — 打包导出
# ============================================================
@app.route('/api/session/<sid>/package', methods=['POST'])
def session_package(sid):
    sess = _get_session(sid)
    if not sess:
        return jsonify(success=False, message='会话不存在')

    data = request.json or {}
    fmt = data.get('format', 'pdf')
    files = data.get('files', [])

    if not files:
        return jsonify(success=False, message='没有图片可打包')

    cache_dir = sess['cache_dir']
    pkg_dir = sess['pkg_dir']
    os.makedirs(pkg_dir, exist_ok=True)

    paths = []
    for f in files:
        p = os.path.join(cache_dir, f)
        if os.path.exists(p):
            paths.append(p)
    if not paths:
        return jsonify(success=False, message='图片文件不存在')

    with sess['lock']:
        vname = Path(sess.get('video_path', '') or 'slides').stem or 'slides'

    try:
        if fmt == 'pdf':
            out = os.path.join(pkg_dir, f'{vname}_整理版.pdf')
            imgs = [Image.open(p).convert('RGB') for p in paths]
            imgs[0].save(out, save_all=True, append_images=imgs[1:])

        elif fmt == 'pptx':
            if not HAS_PPTX:
                return jsonify(success=False, message='未安装 python-pptx，请执行 pip install python-pptx')
            out = os.path.join(pkg_dir, f'{vname}_整理版.pptx')
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            for p in paths:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(out)

        elif fmt == 'zip':
            out = os.path.join(pkg_dir, f'{vname}_整理版.zip')
            with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zf:
                for i, p in enumerate(paths):
                    zf.write(p, f'slide_{i + 1:03d}{Path(p).suffix}')
        else:
            return jsonify(success=False, message=f'不支持的格式: {fmt}')

        return jsonify(success=True, filename=os.path.basename(out))
    except PermissionError:
        return jsonify(success=False,
                       message='文件写入权限被拒绝',
                       hint='请确保目标目录未被占用，或尝试关闭正在使用导出文件的程序。')
    except OSError as e:
        if 'No space' in str(e) or 'disk' in str(e).lower():
            return jsonify(success=False,
                           message='磁盘空间不足，无法导出文件',
                           hint='请清理磁盘空间后重试。')
        return jsonify(success=False, message=f'文件系统错误: {str(e)}',
                       hint='请检查磁盘状态后重试。')
    except Exception as e:
        return jsonify(success=False, message=str(e),
                       hint='导出失败，请重试或换一种导出格式。如果持续出错，请提交 Issue。')


@app.route('/api/session/<sid>/download/<path:filename>')
def session_download(sid, filename):
    sess = _get_session(sid)
    if not sess:
        return jsonify(success=False, message='会话不存在'), 404
    return send_from_directory(sess['pkg_dir'], filename, as_attachment=True)


# ============================================================
#  路由 — 清理单个会话缓存
# ============================================================
@app.route('/api/session/<sid>/cleanup', methods=['POST'])
def session_cleanup(sid):
    sess = _get_session(sid)
    if not sess:
        return jsonify(success=False, message='会话不存在')
    for d in [sess['cache_dir'], sess['pkg_dir']]:
        if os.path.exists(d):
            shutil.rmtree(d)
        os.makedirs(d, exist_ok=True)
    _update_session(sid,
        status='idle', progress=0, message='', saved_count=0,
        video_path='', video_name='', cancel_flag=False,
        eta_seconds=-1, elapsed_seconds=0)
    return jsonify(success=True)


# ============================================================
#  路由 — 全局清理
# ============================================================
@app.route('/api/cleanup-all', methods=['POST'])
def cleanup_all():
    with _sessions_lock:
        sids = list(_sessions.keys())
    for sid in sids:
        sess = _get_session(sid)
        if sess:
            with sess['lock']:
                sess['cancel_flag'] = True
    time.sleep(0.3)
    for sid in sids:
        _delete_session(sid)
    return jsonify(success=True)


# ============================================================
#  路由 — 系统资源监控
# ============================================================
@app.route('/api/system/status')
def system_status():
    result = {
        'cpu_percent': 0,
        'memory_percent': 0,
        'memory_used_gb': 0,
        'memory_total_gb': 0,
        'disk_free_gb': 0,
        'disk_total_gb': 0,
        'disk_percent': 0,
        'active_tasks': _count_running(),
        'total_sessions': len(_sessions),
        'max_sessions': MAX_SESSIONS,
        'warning': None,
        'sessions': _get_all_sessions_summary(),
    }

    if HAS_PSUTIL:
        try:
            result['cpu_percent'] = _cpu_cache['percent']  # 使用后台采样缓存（非阻塞）
            mem = psutil.virtual_memory()
            result['memory_percent'] = mem.percent
            result['memory_used_gb'] = round(mem.used / (1024**3), 1)
            result['memory_total_gb'] = round(mem.total / (1024**3), 1)

            disk = psutil.disk_usage(BASE_DIR)
            result['disk_free_gb'] = round(disk.free / (1024**3), 1)
            result['disk_total_gb'] = round(disk.total / (1024**3), 1)
            result['disk_percent'] = disk.percent
        except Exception as e:
            print(f'[WARN] 获取系统资源信息失败: {e}')

    warnings = []
    if result['cpu_percent'] > CPU_WARN_THRESHOLD:
        warnings.append(f'CPU 使用率过高 ({result["cpu_percent"]:.0f}%)')
    if result['memory_percent'] > MEMORY_WARN_THRESHOLD:
        warnings.append(f'内存使用率过高 ({result["memory_percent"]:.0f}%)')
    if result['disk_free_gb'] < DISK_WARN_THRESHOLD_MB / 1024:
        warnings.append(f'磁盘空间不足 (仅剩 {result["disk_free_gb"]:.1f} GB)')
    if warnings:
        result['warning'] = '；'.join(warnings)

    return jsonify(result)


def _check_resource_warning():
    if not HAS_PSUTIL:
        return None
    try:
        cpu = _cpu_cache['percent']  # 使用后台采样缓存（非阻塞）
        mem = psutil.virtual_memory()
        disk = psutil.disk_usage(BASE_DIR)
        warnings = []
        if cpu > CPU_WARN_THRESHOLD:
            warnings.append(f'CPU 使用率 {cpu:.0f}% 超过 {CPU_WARN_THRESHOLD}%')
        if mem.percent > MEMORY_WARN_THRESHOLD:
            warnings.append(f'内存使用率 {mem.percent:.0f}% 超过 {MEMORY_WARN_THRESHOLD}%')
        if disk.free < DISK_WARN_THRESHOLD_MB * 1024 * 1024:
            warnings.append(f'磁盘空间仅剩 {disk.free / (1024**3):.1f} GB')
        return '；'.join(warnings) if warnings else None
    except Exception:
        return None


# ============================================================
#  路由 — 心跳 & 关闭
# ============================================================
@app.route('/api/heartbeat', methods=['POST'])
def heartbeat():
    global _last_heartbeat, _heartbeat_received
    _last_heartbeat = time.time()
    _heartbeat_received = True
    return jsonify(ok=True)


@app.route('/api/shutdown', methods=['POST'])
def shutdown():
    _do_cleanup()
    print('\n  Shutdown requested, exiting...')
    threading.Timer(0.5, lambda: os._exit(0)).start()
    return jsonify(ok=True)


def _do_cleanup():
    if os.path.exists(SESSIONS_ROOT):
        shutil.rmtree(SESSIONS_ROOT, ignore_errors=True)


def _heartbeat_watcher():
    while True:
        time.sleep(5)
        if not _heartbeat_received:
            continue
        elapsed = time.time() - _last_heartbeat
        if elapsed > HEARTBEAT_TIMEOUT:
            print(f'\n  Browser disconnected for {int(elapsed)}s, shutting down...')
            _do_cleanup()
            print('  Temp files cleaned. Goodbye!')
            time.sleep(0.5)
            os._exit(0)


# ============================================================
#  启动
# ============================================================
def _find_free_port(start=5873):
    for port in range(start, start + 100):
        try:
            s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            s.bind(('127.0.0.1', port))
            s.close()
            return port
        except OSError:
            continue
    return start


if __name__ == '__main__':
    try:
        os.makedirs(SESSIONS_ROOT, exist_ok=True)
        port = _find_free_port(5873)
        url = f'http://127.0.0.1:{port}'

        threading.Timer(1.5, lambda: webbrowser.open(url)).start()

        watcher = threading.Thread(target=_heartbeat_watcher, daemon=True)
        watcher.start()

        print()
        print('=' * 55)
        print('  影幻智提 (VidSlide) v0.2.1 - 多任务版')
        print(f'  浏览器将自动打开: {url}')
        print(f'  临时文件目录: {SESSIONS_ROOT}')
        print(f'  最大并行标签页: {MAX_SESSIONS}')
        print('  关闭浏览器标签页后服务将在 20 秒内自动退出')
        print('  也可以按 Ctrl+C 手动停止')
        print('=' * 55)
        print()

        import atexit
        atexit.register(_do_cleanup)

        app.run(host='127.0.0.1', port=port, debug=False, threaded=True)

    except Exception as e:
        error_detail = traceback.format_exc()
        print(f"！！！发生严重错误！！！\n{error_detail}")
        if sys.stdin is None or not sys.stdout.isatty():
            try:
                import ctypes
                ctypes.windll.user32.MessageBoxW(
                    0,
                    f"影幻智提启动失败！\n\n"
                    f"错误信息：\n{error_detail}\n\n"
                    f"💡 建议操作：\n"
                    f"1. 截图此对话框\n"
                    f"2. 前往 https://github.com/PWO-CHINA/VidSlide/issues 提交 Issue\n"
                    f"3. 在 Issue 中粘贴截图，开发者会尽快修复\n\n"
                    f"常见原因：端口被占用、依赖缺失、杀毒软件拦截",
                    "影幻智提 (VidSlide) - 启动失败",
                    0x10
                )
            except Exception:
                pass
        else:
            print("\n" + "=" * 55)
            print("  💡 建议操作：")
            print("  1. 截图以上错误信息")
            print("  2. 前往 https://github.com/PWO-CHINA/VidSlide/issues 提交 Issue")
            print("  3. 在 Issue 中粘贴截图")
            print("=" * 55)
            input("\n按回车键退出...")
