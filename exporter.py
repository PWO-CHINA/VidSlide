"""
影幻智提 (VidSlide) - 打包导出模块
==================================
负责将提取的幻灯片打包为 PDF、PPTX 或 ZIP 格式。
支持进度回调，可在后台线程中安全调用。

作者: PWO-CHINA
版本: v0.3.2
"""

import os
import zipfile
from pathlib import Path
from PIL import Image

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def package_pdf(paths, output_path, on_progress=None):
    """
    将图片列表打包为 PDF 文件。

    Args:
        paths:       图片文件路径列表
        output_path: 输出 PDF 文件路径
        on_progress: 进度回调 (progress_pct, message)
    """
    total = len(paths)
    if on_progress:
        on_progress(0, '正在生成 PDF…')

    imgs = []
    for i, p in enumerate(paths):
        imgs.append(Image.open(p).convert('RGB'))
        if on_progress:
            on_progress(int((i + 1) / total * 80), f'正在处理第 {i + 1}/{total} 张图片…')

    if on_progress:
        on_progress(85, '正在写入 PDF 文件…')

    imgs[0].save(output_path, save_all=True, append_images=imgs[1:])

    if on_progress:
        on_progress(100, 'PDF 生成完成')


def package_pptx(paths, output_path, on_progress=None):
    """
    将图片列表打包为 PPTX 演示文稿。

    Args:
        paths:       图片文件路径列表
        output_path: 输出 PPTX 文件路径
        on_progress: 进度回调 (progress_pct, message)

    Raises:
        RuntimeError: 如果 python-pptx 未安装
    """
    if not HAS_PPTX:
        raise RuntimeError('未安装 python-pptx，请执行 pip install python-pptx')

    total = len(paths)
    if on_progress:
        on_progress(0, '正在生成 PPTX…')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, p in enumerate(paths):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
        if on_progress:
            on_progress(int((i + 1) / total * 85), f'正在处理第 {i + 1}/{total} 张幻灯片…')

    if on_progress:
        on_progress(90, '正在写入 PPTX 文件…')

    prs.save(output_path)

    if on_progress:
        on_progress(100, 'PPTX 生成完成')


def package_zip(paths, output_path, on_progress=None):
    """
    将图片列表打包为 ZIP 压缩文件。

    Args:
        paths:       图片文件路径列表
        output_path: 输出 ZIP 文件路径
        on_progress: 进度回调 (progress_pct, message)
    """
    total = len(paths)
    if on_progress:
        on_progress(0, '正在创建 ZIP…')

    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for i, p in enumerate(paths):
            zf.write(p, f'slide_{i + 1:03d}{Path(p).suffix}')
            if on_progress:
                on_progress(int((i + 1) / total * 95), f'正在压缩第 {i + 1}/{total} 张图片…')

    if on_progress:
        on_progress(100, 'ZIP 打包完成')


def package_images(paths, pkg_dir, fmt, video_name, on_progress=None):
    """
    主打包入口函数。

    Args:
        paths:       图片文件路径列表
        pkg_dir:     输出目录
        fmt:         格式 ('pdf', 'pptx', 'zip')
        video_name:  视频名称（用于输出文件名）
        on_progress: 进度回调 (progress_pct, message)

    Returns:
        输出文件名（不含路径）

    Raises:
        ValueError:      不支持的格式
        RuntimeError:     缺少依赖
        PermissionError:  文件写入权限被拒绝
        OSError:          磁盘空间不足等文件系统错误
    """
    os.makedirs(pkg_dir, exist_ok=True)

    if fmt == 'pdf':
        out = os.path.join(pkg_dir, f'{video_name}_整理版.pdf')
        package_pdf(paths, out, on_progress)
    elif fmt == 'pptx':
        out = os.path.join(pkg_dir, f'{video_name}_整理版.pptx')
        package_pptx(paths, out, on_progress)
    elif fmt == 'zip':
        out = os.path.join(pkg_dir, f'{video_name}_整理版.zip')
        package_zip(paths, out, on_progress)
    else:
        raise ValueError(f'不支持的格式: {fmt}')

    return os.path.basename(out)
